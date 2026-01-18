# ============================================================================
# PowerPoint 匯出模組 (PowerPoint Export Module)
# ============================================================================
# 說明：將簡報資料匯出為 PPTX 格式檔案
# 功能：使用 python-pptx 函式庫生成專業簡報
# 特色：捷運站牌美學設計，16:9 寬螢幕格式
# 架構：採用建造者模式 (Builder Pattern)
# ============================================================================

# 標準函式庫匯入
import base64  # Base64 編碼解碼功能
import io  # 記憶體串流處理
from typing import List  # 型別提示

# 第三方函式庫匯入 - python-pptx
from pptx import Presentation  # PowerPoint 簡報物件
from pptx.util import Inches, Pt  # 單位轉換工具（英吋、點）
from pptx.dml.color import RGBColor  # RGB 色彩物件
from pptx.enum.shapes import MSO_SHAPE  # 形狀類型列舉
from pptx.enum.text import PP_ALIGN  # 文字對齊方式列舉

# 專案內部模組匯入
from src.models import SlideContent  # 簡報頁面內容資料模型
from src.styles import StyleProfile  # 風格設定檔資料模型

# ============================================================================
# PowerPoint 匯出器類別 (PPTX Exporter)
# ============================================================================
class PptxExporter:
    """
    PowerPoint 簡報匯出器
    
    功能：
    - 將結構化簡報資料轉換為 PPTX 檔案
    - 支援背景圖片嵌入
    - 實作簡報視覺風格
    - 提供進度指示器與配色系統
    
    設計特色：
    - 16:9 寬螢幕比例
    - 簡報路線色彩系統
    - 站點編號徽章
    - 玻璃擬態文字背景
    - 右上角進度點指示器
    
    使用範例：
        exporter = PptxExporter()
        exporter.export(slides, "output.pptx", style_profile)
    """
    
    # ------------------------------------------------------------------------
    # 初始化方法 (Constructor)
    # ------------------------------------------------------------------------
    def __init__(self):
        """
        初始化 PowerPoint 匯出器
        
        功能：
        - 設定簡報尺寸為 16:9 寬螢幕格式
        
        尺寸標準：
        - 寬度：13.333 英吋 (16 單位)
        - 高度：7.5 英吋 (9 單位)
        - 比例：16:9
        """
        # 設定簡報寬度（16:9 寬螢幕格式）
        self.width = Inches(13.333)
        
        # 設定簡報高度
        self.height = Inches(7.5)
    
    # ------------------------------------------------------------------------
    # 公開方法：匯出簡報 (Public Method: Export Presentation)
    # ------------------------------------------------------------------------
    def export(self, slides: List[SlideContent], output_path: str, style: StyleProfile):
        """
        匯出簡報為 PPTX 檔案
        
        功能：
        - 建立新的 PowerPoint 簡報
        - 遍歷所有簡報頁面並建立投影片
        - 應用視覺風格與配色系統
        - 儲存為 PPTX 檔案
        
        參數：
            slides (List[SlideContent]): 簡報頁面列表
            output_path (str): 輸出檔案路徑（例如："presentation.pptx"）
            style (StyleProfile): 風格設定檔（包含配色映射）
        
        處理流程：
        1. 建立新簡報並設定尺寸
        2. 遍歷每一頁簡報內容
        3. 為每頁添加背景、UI 元素、文字
        4. 儲存檔案
        """
        # 步驟 1: 建立新的 PowerPoint 簡報物件
        prs = Presentation()
        
        # 步驟 2: 設定簡報尺寸為 16:9
        prs.slide_width = self.width
        prs.slide_height = self.height

        # 步驟 3: 遍歷所有簡報頁面
        for idx, content in enumerate(slides):
            # 步驟 3.1: 建立空白投影片
            # prs.slide_layouts[6] 是完全空白的版型
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # 步驟 3.2: 設定背景圖片（如果存在）
            if content.background_image_base64:
                self._set_background(slide, content.background_image_base64)

            # 步驟 3.3: 取得當前頁面的路線顏色
            # 根據框架章節從風格配色映射中取得對應的顏色
            hex_color = style.logic_color_map.get(content.framework_section, "#333333")
            
            # 將十六進位色碼轉換為 RGB 物件
            rgb_color = self._hex_to_rgb(hex_color)

            # 步驟 3.4: 繪製捷運風格 UI 元素
            # 包含：路線色條、站點徽章、進度指示器
            self._draw_presentation_ui(slide, idx + 1, rgb_color, len(slides))

            # 步驟 3.5: 繪製文字內容區域
            self._draw_text_content(slide, content, rgb_color)

        # 步驟 4: 儲存簡報檔案
        prs.save(output_path)
    
    # ------------------------------------------------------------------------
    # 私有方法：繪製文字內容 (Private Method: Draw Text Content)
    # ------------------------------------------------------------------------
    def _draw_text_content(self, slide, content: SlideContent, rgb_color: RGBColor):
        """
        繪製簡報頁面的文字內容區域
        
        功能：
        - 建立半透明白色背景色塊（玻璃擬態效果）
        - 添加標題與內文文字
        - 添加左側彩色裝飾條
        
        參數：
            slide: PowerPoint 投影片物件
            content (SlideContent): 簡報頁面內容
            rgb_color (RGBColor): 路線顏色（用於裝飾條）
        
        設計理念：
        - 玻璃擬態設計：半透明白色背景
        - 視覺層次：標題與內文分離
        - 色彩強調：左側裝飾條對應路線顏色
        """
        # ---- 標題區域 ----
        
        # 步驟 1: 建立標題背景色塊（圓角矩形）
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,  # 圓角矩形
            Inches(1.2),  # X 位置
            Inches(1.65),  # Y 位置
            Inches(9),    # 寬度
            Inches(1.1)   # 高度
        )
        
        # 設定標題背景為白色，96% 不透明（玻璃擬態效果）
        title_bg.fill.solid()  # 實心填充
        title_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白色
        title_bg.fill.transparency = 0.04  # 透明度 4% (96% 不透明)
        title_bg.line.fill.background()  # 移除邊框
        title_bg.shadow.inherit = False  # 不繼承陰影
        
        # 步驟 2: 建立內文背景色塊（圓角矩形）
        body_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,  # 圓角矩形
            Inches(1.2),  # X 位置
            Inches(2.95),  # Y 位置
            Inches(9),     # 寬度
            Inches(3.2)    # 高度
        )
        
        # 設定內文背景為白色，94% 不透明
        body_bg.fill.solid()  # 實心填充
        body_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白色
        body_bg.fill.transparency = 0.06  # 透明度 6% (94% 不透明)
        body_bg.line.fill.background()  # 移除邊框
        body_bg.shadow.inherit = False  # 不繼承陰影
        
        # 步驟 3: 添加標題文字方塊
        title_box = slide.shapes.add_textbox(
            Inches(1.4),   # X 位置（比背景稍微內縮）
            Inches(1.75),  # Y 位置
            Inches(8.6),   # 寬度
            Inches(0.9)    # 高度
        )
        
        # 設定標題文字樣式（44pt 粗體深灰色）
        self._add_text(
            title_box,
            content.title,
            size=44,
            bold=True,
            color=RGBColor(26, 26, 26)  # 深灰色
        )
        
        # 步驟 4: 添加左側彩色裝飾條
        # 用途：視覺強調，對應路線顏色
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,  # 矩形
            Inches(1.2),   # X 位置（與背景對齊）
            Inches(1.65),  # Y 位置
            Inches(0.08),  # 寬度（細長條）
            Inches(1.1)    # 高度（與標題背景相同）
        )
        
        # 設定裝飾條顏色為路線顏色
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = rgb_color
        accent_bar.line.fill.background()  # 移除邊框
        
        # 步驟 5: 添加內文文字方塊
        body_box = slide.shapes.add_textbox(
            Inches(1.5),   # X 位置
            Inches(3.15),  # Y 位置
            Inches(8.3),   # 寬度
            Inches(2.8)    # 高度
        )
        
        # 設定內文文字樣式（24pt 一般字重深灰色）
        self._add_text(
            body_box,
            content.body_text,
            size=24,
            bold=False,
            color=RGBColor(42, 42, 42)  # 中深灰色
        )
    
    # ------------------------------------------------------------------------
    # 私有方法：繪製簡報視覺系統的 UI 元素 (Private Method: Draw Presentation UI)
    # ------------------------------------------------------------------------
    def _draw_presentation_ui(self, slide, index: int, rgb_color: RGBColor, total_slides: int):
        """
        繪製簡報視覺系統的 UI 元素
        
        功能：
        - 繪製橫向路線色條
        - 繪製站點編號徽章
        - 繪製右上角進度指示器
        
        參數：
            slide: PowerPoint 投影片物件
            index (int): 當前頁碼（從 1 開始）
            rgb_color (RGBColor): 路線顏色
            total_slides (int): 總頁數
        
        設計元素：
        A. 路線色條：橫跨整個投影片的彩色橫條
        B. 站點徽章：左上角的圓形站點編號
        C. 進度指示器：右上角的點狀進度條
        """
        # ---- A. 路線色條 (Metro Line) ----
        
        # 步驟 A1: 建立橫向色條
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,  # 矩形
            0,                    # X 位置（從最左側開始）
            Inches(1.2),         # Y 位置
            self.width,          # 寬度（整個投影片寬度）
            Inches(0.18)         # 高度（細長條）
        )
        
        # 步驟 A2: 設定色條樣式
        line.fill.solid()  # 實心填充
        line.fill.fore_color.rgb = rgb_color  # 使用路線顏色
        line.line.fill.background()  # 移除邊框
        line.shadow.inherit = False  # 不繼承陰影

        # ---- B. 站點徽章 (Station Badge) ----
        
        # 步驟 B1: 建立圓形徽章
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,    # 橢圓形（正圓）
            Inches(0.5),       # X 位置
            Inches(0.75),      # Y 位置
            Inches(0.9),       # 寬度
            Inches(0.9)        # 高度（寬高相同 = 正圓）
        )
        
        # 步驟 B2: 設定圓形樣式（白底 + 彩色粗邊框）
        circle.fill.solid()  # 實心填充
        circle.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白色填充
        circle.line.color.rgb = rgb_color  # 邊框使用路線顏色
        circle.line.width = Pt(6)  # 邊框寬度 6 點（粗邊框）
        circle.shadow.inherit = False  # 不繼承陰影

        # 步驟 B3: 在圓形中添加站點編號文字
        tf = circle.text_frame  # 取得文字框架
        p = tf.paragraphs[0]    # 取得第一個段落
        
        # 格式化編號為兩位數（例如：01, 02, 10, 11）
        p.text = str(index).zfill(2)
        
        # 設定文字樣式
        p.font.size = Pt(28)  # 字體大小 28 點
        p.font.bold = True    # 粗體
        p.font.color.rgb = rgb_color  # 文字顏色使用路線顏色
        p.alignment = PP_ALIGN.CENTER  # 水平置中
        
        # 設定文字框架屬性
        tf.margin_top = Inches(0.18)  # 上邊距（微調垂直位置）
        tf.vertical_anchor = 1  # 垂直置中對齊
        
        # ---- C. 進度指示器 (Progress Dots) ----
        
        # 步驟 C1: 計算進度點的尺寸與間距
        dot_size = Inches(0.15)  # 每個點的直徑
        dot_gap = Inches(0.12)   # 點與點之間的間距
        
        # 步驟 C2: 計算起始位置（從右側往左排列）
        start_x = self.width - Inches(0.5) - (total_slides * (dot_size + dot_gap))
        start_y = Inches(0.35)  # 距離頂部的位置
        
        # 步驟 C3: 遍歷所有頁面，繪製進度點
        for i in range(1, total_slides + 1):
            # 計算當前點的 X 位置
            dot_x = start_x + (i - 1) * (dot_size + dot_gap)
            
            # 建立圓形進度點
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,  # 橢圓形（正圓）
                dot_x,           # X 位置
                start_y,         # Y 位置
                dot_size,        # 寬度
                dot_size         # 高度（正圓）
            )
            
            # 設定進度點樣式
            dot.fill.solid()  # 實心填充
            
            if i == index:
                # 當前頁面：亮白色 + 彩色邊框（高亮顯示）
                dot.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白色填充
                dot.line.color.rgb = rgb_color  # 彩色邊框
                dot.line.width = Pt(2.5)  # 邊框寬度
            else:
                # 其他頁面：半透明灰色（淡化顯示）
                dot.fill.fore_color.rgb = RGBColor(200, 200, 200)  # 淺灰色
                dot.fill.transparency = 0.5  # 透明度 50%
                dot.line.fill.background()  # 移除邊框
            
            # 不繼承陰影
            dot.shadow.inherit = False
    
    # ------------------------------------------------------------------------
    # 私有方法：設定背景圖片 (Private Method: Set Background)
    # ------------------------------------------------------------------------
    def _set_background(self, slide, b64_str: str):
        """
        設定投影片的背景圖片
        
        功能：
        - 將 Base64 編碼的圖片解碼
        - 插入為投影片背景（填滿整個投影片）
        - 提供錯誤處理機制
        
        參數：
            slide: PowerPoint 投影片物件
            b64_str (str): Base64 編碼的圖片字串
        
        技術細節：
        - 支援 data URI 格式（自動移除前綴）
        - 使用記憶體串流避免建立臨時檔案
        - 圖片會被置於最底層（Z-order）
        
        注意：
        - 第一個添加的形狀自然在最底層
        - 後續添加的元素會疊在圖片上方
        """
        try:
            # 步驟 1: 處理 Base64 字串格式
            # 如果包含 data URI 前綴（例如："data:image/png;base64,..."），則移除
            if "," in b64_str:
                b64_str = b64_str.split(",")[1]
            
            # 步驟 2: 解碼 Base64 字串為二進位資料
            img_data = base64.b64decode(b64_str)
            
            # 步驟 3: 使用記憶體串流插入圖片
            # BytesIO: 將二進位資料包裝為類檔案物件
            # 參數：X=0, Y=0, Width=投影片寬, Height=投影片高（填滿整個投影片）
            slide.shapes.add_picture(
                io.BytesIO(img_data),  # 圖片資料串流
                0,                      # X 位置（左上角）
                0,                      # Y 位置（左上角）
                self.width,            # 寬度（填滿）
                self.height            # 高度（填滿）
            )
            
            # 技術說明：
            # 由於這是第一個添加的形狀，它會自動位於最底層
            # 後續添加的所有元素（文字、UI）都會在圖片上方
            
        except Exception:
            # 錯誤處理：如果圖片解碼或插入失敗，靜默忽略
            # 原因：避免單張圖片錯誤導致整個簡報匯出失敗
            pass
    
    # ------------------------------------------------------------------------
    # 私有方法：添加文字 (Private Method: Add Text)
    # ------------------------------------------------------------------------
    def _add_text(self, shape, text: str, size: int, bold: bool, color: RGBColor):
        """
        為形狀添加格式化文字
        
        功能：
        - 設定文字內容與樣式
        - 支援自動換行
        - 統一字體為微軟正黑體
        
        參數：
            shape: PowerPoint 形狀物件（包含 text_frame）
            text (str): 文字內容
            size (int): 字體大小（點 pt）
            bold (bool): 是否粗體
            color (RGBColor): 文字顏色
        
        設定項目：
        - 字體：Microsoft JhengHei（微軟正黑體）
        - 自動換行：啟用
        - 字體大小、粗細、顏色：依參數設定
        """
        # 步驟 1: 取得文字框架
        tf = shape.text_frame
        
        # 步驟 2: 啟用自動換行
        tf.word_wrap = True
        
        # 步驟 3: 取得第一個段落
        p = tf.paragraphs[0]
        
        # 步驟 4: 設定文字內容
        p.text = text
        
        # 步驟 5: 設定字體屬性
        p.font.name = "Microsoft JhengHei"  # 微軟正黑體（適合中文）
        p.font.size = Pt(size)              # 字體大小
        p.font.bold = bold                  # 粗體設定
        p.font.color.rgb = color            # 文字顏色
    
    # ------------------------------------------------------------------------
    # 私有方法：十六進位轉 RGB (Private Method: Hex to RGB)
    # ------------------------------------------------------------------------
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """
        將十六進位色碼轉換為 RGBColor 物件
        
        功能：
        - 解析十六進位色碼（例如："#E3002C"）
        - 轉換為 PowerPoint 使用的 RGBColor 物件
        
        參數：
            hex_color (str): 十六進位色碼（例如："#E3002C" 或 "E3002C"）
        
        回傳：
            RGBColor: PowerPoint RGB 顏色物件
        
        轉換邏輯：
        1. 移除 # 符號（如果存在）
        2. 將十六進位字串拆分為 R、G、B 三個部分
        3. 每部分轉換為 0-255 的整數
        4. 建立 RGBColor 物件
        
        範例：
            "#E3002C" → RGBColor(227, 0, 44)
        """
        # 步驟 1: 移除 # 符號（如果存在）
        h = hex_color.lstrip('#')
        
        # 步驟 2: 解析並轉換
        # h[0:2]: 前兩個字元（R）
        # h[2:4]: 中間兩個字元（G）
        # h[4:6]: 後兩個字元（B）
        # int(x, 16): 將十六進位字串轉為整數
        return RGBColor(*tuple(int(h[i:i+2], 16) for i in (0, 2, 4)))

# ============================================================================
# 模組定義完成
# ============================================================================
#
# 技術細節：
# 1. python-pptx：Python 的 PowerPoint 操作函式庫
# 2. Z-order：形狀的堆疊順序（先添加的在底層）
# 3. 玻璃擬態：半透明白色背景的設計風格
# 4. 單位轉換：Inches（英吋）、Pt（點）
#
# 擴充指南：
# 1. 新增模板：建立不同的 _draw_* 方法實作新風格
# 2. 自訂動畫：使用 python-pptx-interface 添加動畫效果
# 3. 圖表支援：整合 python-pptx 的圖表功能
#
# 設計模式：
# - 建造者模式 (Builder Pattern)：逐步建構複雜的簡報物件
# - 策略模式 (Strategy Pattern)：可切換不同的視覺風格
# ============================================================================